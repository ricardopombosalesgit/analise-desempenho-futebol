import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from datetime import datetime
from io import BytesIO
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os
from openpyxl import Workbook
from openpyxl.drawing.image import Image

# Configuração da página
st.set_page_config(page_title="Futebol: Análise Individual", layout="wide")

# Título da aplicação
st.title("Futebol: Análise Individual")
st.write("Feito por Ricardo Pombo Sales")

# Criando o seletor de posição de jogo
menu = st.selectbox(
    "Selecione a Posição de Jogo:",
    ["Selecione uma Posição", "Goleiro(a)", "Zagueiro(a)", "Lateral", "Volante", "Meia", "Extremo(a)", "Atacante", "Geral"]
)

# Função para salvar a imagem de um gráfico como arquivo temporário
def salvar_imagem(fig, filename):
    fig.write_image(filename)
    return filename

# Função para gerar relatório em Excel com gráficos
def gerar_excel(data, imagens):
    # Criação do workbook e planilha
    wb = Workbook()
    ws = wb.active
    ws.title = "Avaliação"

    # Adicionando os dados na planilha
    for i, (key, value) in enumerate(data.items(), start=1):
        ws.cell(row=i, column=1, value=key)
        ws.cell(row=i, column=2, value=value)

    # Adicionando os gráficos na planilha
    for i, img_path in enumerate(imagens, start=1):
        img = Image(img_path)
        # Define a posição onde a imagem será inserida
        ws.add_image(img, f"E{i * 15}")  # Coloca cada gráfico a partir de uma linha nova (ajuste conforme necessário)

    # Salvando o Excel em memória
    excel_file = BytesIO()
    wb.save(excel_file)
    excel_file.seek(0)
    return excel_file

# Função para gerar relatório em PDF com gráficos
def gerar_pdf(data, imagens):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Relatório de Avaliação", ln=True, align='C')
    for key, value in data.items():
        pdf.cell(200, 10, txt=f"{key}: {value}", ln=True)

    # Adicionando imagens no PDF
    for imagem in imagens:
        pdf.add_page()  # Adiciona uma nova página para cada imagem
        pdf.image(imagem, x=10, y=20, w=180)  # Ajusta a posição e o tamanho da imagem no PDF

    # Salva o PDF como string de bytes
    pdf_file = BytesIO()
    pdf_content = pdf.output(dest='S').encode('latin1')  # 'S' retorna o PDF como string de bytes
    pdf_file.write(pdf_content)  # Escreve o conteúdo no BytesIO
    pdf_file.seek(0)  # Posiciona o cursor no início do BytesIO
    return pdf_file

# Função para gerar relatório em DOCX (Word) com gráficos
def gerar_docx(data, imagens):
    doc = Document()
    doc.add_heading('Relatório de Avaliação', 0)
    for key, value in data.items():
        doc.add_paragraph(f"{key}: {value}")
    for imagem in imagens:
        doc.add_picture(imagem, width=Inches(6))
    docx_file = BytesIO()
    doc.save(docx_file)
    docx_file.seek(0)
    return docx_file

# Função para gerar relatório em PPTX (PowerPoint) com gráficos
def gerar_pptx(data, imagens):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Relatório de Avaliação"
    content.text = "\n".join([f"{key}: {value}" for key, value in data.items()])

    for imagem in imagens:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Gráfico"
        slide.shapes.add_picture(imagem, Inches(1), Inches(1), width=Inches(8))

    pptx_file = BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    return pptx_file

# Função para exibir campos de entrada de dados e avaliações com base na seleção
def show_page(option):
    # Função comum para entrada de dados
    def entrada_dados():
        col1, col2 = st.columns(2)
        
        nome = col1.text_input("Nome:")
        data_nascimento = col1.text_input("Data de Nascimento (dd/mm/aaaa):", value="DD/MM/AAAA")
        temporada = col1.text_input("Temporada:")
        altura = col1.number_input("Altura (cm):", 0, 230, 0)
        selecao_nacional = col1.selectbox("Convocação em Seleção Nacional?", ["Sim", "Não"])
        categoria = col1.text_input("Categoria:")
        minutos_totais = col1.number_input("Minutos Totais na Temporada:", 0, 100000, 0)

        data_analise = col2.date_input("Data da Análise (dd/mm/aaaa):", value=datetime.today())
        idade = col2.number_input("Idade:", min_value=0, max_value=100, value=0, step=1)
        nacionalidade = col2.text_input("Nacionalidade:")
        peso = col2.number_input("Peso (kg):", 0, 200, 0)
        clube = col2.text_input("Clube:")
        num_jogos = col2.number_input("Número de Jogos na Temporada:", 0, 1000, 0)
        pe_dominante = col2.selectbox("Pé Dominante:", ["Direito", "Esquerdo", "Ambidestro"])
        cartoes_amarelos = col1.number_input("Cartões Amarelos:", 0, 50, 0)
        cartoes_vermelhos = col2.number_input("Cartões Vermelhos:", 0, 20, 0)

        data = {
            'Nome': nome, 'Nacionalidade': nacionalidade, 'Data de Nascimento': data_nascimento,
            'Idade': idade, 'Altura': altura, 'Peso': peso, 'Seleção Nacional': selecao_nacional, 'Clube': clube,
            'Categoria': categoria, 'Data da Análise': data_analise.strftime('%d/%m/%Y'), 'Temporada': temporada,
            'Número de Jogos': num_jogos, 'Minutos Totais': minutos_totais, 'Pé Dominante': pe_dominante,
            'Cartões Amarelos': cartoes_amarelos, 'Cartões Vermelhos': cartoes_vermelhos
        }
        return data

    # Função para avaliações específicas do goleiro(a)
    def avaliacoes_goleiro():
        st.subheader("Avaliação de Ações Ofensivas do Goleiro(a)")
        acoes_ofensivas = {
            'Percepção do jogo ofensivo': st.slider('Percepção do jogo ofensivo', 5, 10, 7),
            'Passe com pé': st.slider('Passe com pé', 5, 10, 7),
            'Lançamento longo com pé': st.slider('Lançamento longo com pé', 5, 10, 7),
            'Lançamento curto com as mãos': st.slider('Lançamento curto com as mãos', 5, 10, 7),
            'Lançamento longo com as mãos': st.slider('Lançamento longo com as mãos', 5, 10, 7),
            'Tiro de meta curto': st.slider('Tiro de meta curto', 5, 10, 7),
            'Tiro de meta longo': st.slider('Tiro de meta longo', 5, 10, 7)
        }

        st.subheader("Avaliação de Ações Defensivas do Goleiro(a)")
        acoes_defensivas = {
            'Percepção do jogo defensivo': st.slider('Percepção do jogo defensivo', 5, 10, 7),
            'Controle do jogo defensivo': st.slider('Controle do jogo defensivo', 5, 10, 7),
            'Comunicação': st.slider('Comunicação', 5, 10, 7),
            'Vigor nas jogadas aéreas': st.slider('Vigor nas jogadas aéreas', 5, 10, 7),
            'Proteção da profundidade': st.slider('Proteção da profundidade', 5, 10, 7),
            'Encaixe': st.slider('Encaixe', 5, 10, 7),
            'Espalmada': st.slider('Espalmada', 5, 10, 7),
            'Interceptação de Cruzamentos': st.slider('Interceptação de Cruzamentos', 5, 10, 7),
            'Defesa com os pés': st.slider('Defesa com os pés', 5, 10, 7),
            'Defesa com as mãos': st.slider('Defesa com as mãos', 5, 10, 7),
            'Defesa com outra parte do corpo': st.slider('Defesa com outra parte do corpo', 5, 10, 7),
            'Intervenção': st.slider('Intervenção', 5, 10, 7)
        }

        return acoes_ofensivas, acoes_defensivas

    # Função comum para avaliações de outras posições
    def avaliacoes(acoes_adicionais=None):
        st.subheader("Avaliação de Ações Ofensivas")
        acoes_ofensivas = {
            'Gols Marcados': st.number_input('Gols Marcados:', 0, 50, 0),
            'Assistências': st.number_input('Assistências:', 0, 50, 0),
            'Percepção do jogo ofensivo': st.slider('Percepção do jogo ofensivo', 5, 10, 7),
            'Cruzamentos': st.slider('Cruzamentos:', 5, 10, 7),
            'Progressão do jogo': st.slider('Progressão do jogo', 5, 10, 7),
            'Passe geral': st.slider('Passe geral', 5, 10, 7),
            'Passe curto': st.slider('Passe curto', 5, 10, 7),
            'Passe médio': st.slider('Passe médio', 5, 10, 7),
            'Passe longo': st.slider('Passe longo', 5, 10, 7),
            'Passe chave': st.slider('Passe chave', 5, 10, 7),
            'Finalizações': st.slider('Finalizações', 5, 10, 7),
            'Finalizações de cabeça': st.slider('Finalizações de cabeça', 5, 10, 7),
            'Finalizações com o pé': st.slider('Finalizações com o pé', 5, 10, 7),
            'Finalizações curta distância': st.slider('Finalizações curta distância', 5, 10, 7),
            'Finalizações longa distância': st.slider('Finalizações longa distância', 5, 10, 7),
            'Chances de gols': st.slider('Chances de gols', 5, 10, 7),
            'Chances criadas': st.slider('Chances criadas', 5, 10, 7),
            'Desafio no ataque': st.slider('Desafio no ataque', 5, 10, 7),
            'Transição ofensiva': st.slider('Transição ofensiva', 5, 10, 7),
            'Desafios': st.slider('Desafios', 5, 10, 7)
        }

        if acoes_adicionais:
            acoes_ofensivas.update(acoes_adicionais)

        st.subheader("Avaliação de Ações Defensivas")
        acoes_defensivas = {
            'Percepção do jogo defensivo': st.slider('Percepção do jogo defensivo', 5, 10, 7),
            'Controle do jogo defensivo': st.slider('Controle do jogo defensivo', 5, 10, 7),
            'Comunicação': st.slider('Comunicação', 5, 10, 7),
            'Vigor nas jogadas aéreas': st.slider('Vigor nas jogadas aéreas', 5, 10, 7),
            'Proteção da profundidade': st.slider('Proteção da profundidade', 5, 10, 7),
            'Desarmes': st.slider('Desarmes', 5, 10, 7),
            'Disputa': st.slider('Disputa', 5, 10, 7),
            'Perda de bola': st.slider('Perda de bola', 5, 10, 7),
            'Desafio na defesa': st.slider('Desafio na defesa', 5, 10, 7),
            'Falhas em Gols': st.slider('Falhas em Gols', 5, 10, 7),
            'Recuperação de bolas': st.slider('Recuperação de bolas', 5, 10, 7),
            'Interceptações': st.slider('Interceptações', 5, 10, 7),
            'Rebotes': st.slider('Rebotes', 5, 10, 7),
            'Disputas aéreas': st.slider('Disputas aéreas', 5, 10, 7),
            'Transição defensiva': st.slider('Transição defensiva', 5, 10, 7)
        }

        return acoes_ofensivas, acoes_defensivas

    # Função para criar gráficos de avaliação
    def criar_graficos(acoes_ofensivas, acoes_defensivas):
        labels = list(acoes_ofensivas.keys()) + list(acoes_defensivas.keys())
        values = list(acoes_ofensivas.values()) + list(acoes_defensivas.values())

        fig_radar = go.Figure(data=go.Scatterpolar(
            r=values,
            theta=labels,
            fill='toself'
        ))

        fig_radar.update_layout(
            polar=dict(
                radialaxis=dict(
                    visible=True,
                    range=[0, 10]
                )),
            showlegend=False
        )

        st.plotly_chart(fig_radar, use_container_width=True)

        fig_bar = px.bar(
            x=labels,
            y=values,
            title="Comparação de Ações Ofensivas e Defensivas",
            labels={'x': 'Ações', 'y': 'Avaliação'}
        )
        st.plotly_chart(fig_bar, use_container_width=True)

        # Ajuste para garantir que as listas x e y tenham o mesmo comprimento
        min_len = min(len(list(acoes_ofensivas.values())), len(list(acoes_defensivas.values())))
        fig_scatter = px.scatter(
            x=list(acoes_ofensivas.values())[:min_len],
            y=list(acoes_defensivas.values())[:min_len],
            title="Correlação entre Ações Ofensivas e Defensivas",
            labels={'x': 'Ações Ofensivas', 'y': 'Ações Defensivas'}
        )
        st.plotly_chart(fig_scatter, use_container_width=True)

        fig_line = px.line(
            x=[i for i in range(1, len(acoes_ofensivas) + 1)],
            y=list(acoes_ofensivas.values()),
            title="Evolução das Ações Ofensivas ao Longo dos Jogos",
            labels={'x': 'Jogos', 'y': 'Avaliações Ofensivas'}
        )
        st.plotly_chart(fig_line, use_container_width=True)

        return [fig_radar, fig_bar, fig_scatter, fig_line]

    # Mostrar a parte de entrada de dados somente se uma posição válida for selecionada
    if option != "Selecione uma Posição":
        if option == "Goleiro(a)":
            st.subheader("Análise de Goleiro(a)")
            data = entrada_dados()
            acoes_ofensivas, acoes_defensivas = avaliacoes_goleiro()
            graficos = criar_graficos(acoes_ofensivas, acoes_defensivas)
        else:
            st.subheader(f"Análise de {option}")
            data = entrada_dados()
            acoes_ofensivas, acoes_defensivas = avaliacoes()
            graficos = criar_graficos(acoes_ofensivas, acoes_defensivas)

        # Salvando gráficos e criando os arquivos para download
        with tempfile.TemporaryDirectory() as temp_dir:
            imagens = [salvar_imagem(fig, os.path.join(temp_dir, f"grafico_{i}.png")) for i, fig in enumerate(graficos)]

            st.download_button(
                label="Baixar Relatório em Excel",
                data=gerar_excel(data, imagens),
                file_name=f"relatorio_{option.lower()}_{data['Nome']}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

            st.download_button(
                label="Baixar Relatório em PDF",
                data=gerar_pdf(data, imagens),
                file_name=f"relatorio_{option.lower()}_{data['Nome']}.pdf",
                mime="application/pdf"
            )

            st.download_button(
                label="Baixar Relatório em Word",
                data=gerar_docx(data, imagens),
                file_name=f"relatorio_{option.lower()}_{data['Nome']}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

            st.download_button(
                label="Baixar Relatório em PowerPoint",
                data=gerar_pptx(data, imagens),
                file_name=f"relatorio_{option.lower()}_{data['Nome']}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

# Exibindo a página correspondente à seleção
show_page(menu)